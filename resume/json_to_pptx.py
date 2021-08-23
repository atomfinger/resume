from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt, Inches

from resume.json_to_md import format_date
from resume.resume_types import Resume, Work


def convert_to_pptx(resume: Resume, destination: str):
    presentation = Presentation()
    configure_title_slide(presentation.slides.add_slide(presentation.slide_layouts[0]), resume)
    configure_works_experience_slides(presentation, resume)
    add_happiness_slide(presentation.slides.add_slide(presentation.slide_layouts[5]), resume)
    add_education_slides(presentation.slides.add_slide(presentation.slide_layouts[1]), resume)
    add_personality_slide(presentation.slides.add_slide(presentation.slide_layouts[5]))
    add_contact_slide(presentation.slides.add_slide(presentation.slide_layouts[0]), resume)
    presentation.save(destination)


def configure_title_slide(slide, resume: Resume):
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f'{resume.basics.name}'
    subtitle.text = f'{resume.basics.summary}'


def configure_works_experience_slides(presentation, resume: Resume):
    for work in resume.work:
        configure_work_experience_slides(presentation.slides.add_slide(presentation.slide_layouts[1]), work)


def configure_work_experience_slides(slide, work: Work):
    title = slide.shapes.title
    title.text = f'{work.position}@{work.name} ({format_date(work.start_date)}-{format_date(work.end_date, "Current")})'
    title.text_frame.paragraphs[0].font.size = Pt(24)
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for highlight in work.highlights:
        p = tf.add_paragraph()
        p.text = highlight
        p.level = 1


def add_happiness_slide(slide, resume: Resume):
    title = slide.shapes.title
    title.text = f'Happiness rating before and after {resume.basics.name} joining a company'
    title.text_frame.paragraphs[0].font.size = Pt(15)
    add_happiness_chart(slide)


def add_happiness_chart(slide):
    chart_data = ChartData()
    chart_data.categories = ['Q1', 'Q2', f'Hiring John', 'Q3', 'Q4']
    chart_data.add_series('Happiness', (30.0, 23.0, 17.2, 69.69, 120.00))
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = False
    chart.series[0].smooth = True
    chart.value_axis.visible = False
    chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)


def add_education_slides(slide, resume: Resume):
    title = slide.shapes.title
    title.text = "Education"
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for education in resume.education:
        p = tf.add_paragraph()
        p.text = f'{education.area}/{education.institution} ' \
                 f'({format_date(education.start_date)}-{format_date(education.end_date, "Current")})'
        p.font.size = Pt(20)
        p.level = 1


def add_contact_slide(slide, resume:Resume):
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f'Contact him today'
    subtitle.text = f'{resume.basics.email} or {resume.basics.url}'


def add_personality_slide(slide):
    title = slide.shapes.title
    title.text = f'Personality'

    chart_data = ChartData()
    chart_data.categories = ['Nerd', 'Professional', 'Fun', 'Intelligent', 'Swag']
    chart_data.add_series('', (0.50, 0.17, 0.30, 0.01, 0.02))

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0%'
    # data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
